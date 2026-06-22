"""Tạo slide PPTX từ nội dung đã cấu trúc"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pathlib import Path
from typing import Dict, Any, List, Optional
import asyncio
import re

from services.slide_charts import normalize_chart_spec, _DATA_KEYWORDS as _DATA_SLIDE_KEYWORDS
from services.slide_tables import slide_has_table_or_body
from services.content.slide_normalizer import _is_generic_title

# Preset giao diện PPTX (chọn từ demo / API). Khác với theme theo từ khóa title (blue/green/...).
VALID_SLIDE_PRESETS = frozenset({"corporate", "modern", "minimal"})


class SlideGenerator:
    """Tạo file PowerPoint từ nội dung đã cấu trúc"""
    
    def _lighten_color(self, color: RGBColor, amount: int = 50) -> RGBColor:
        """Tạo màu nhạt hơn từ RGBColor"""
        return RGBColor(
            min(255, color[0] + amount),
            min(255, color[1] + amount),
            min(255, color[2] + amount)
        )
    
    def __init__(self):
        # Widescreen 16:9
        self.slide_width = Inches(13.333)
        self.slide_height = Inches(7.5)

        # Layout constants
        self.margin_x = Inches(0.8)
        self.margin_y = Inches(0.6)
        self.title_h = Inches(1.0)
        self.footer_h = Inches(0.35)

        # Fonts - dùng fonts phổ biến trên Windows
        self.font_title = "Segoe UI"
        self.font_body = "Segoe UI"
        
        # Theme colors - Professional themes
        self.themes = {
            "default": {
                "bg_start": RGBColor(241, 245, 249),  # Slate-100
                "bg_end": RGBColor(255, 255, 255),   # White
                "title": RGBColor(30, 64, 175),      # Blue-700
                "body": RGBColor(15, 23, 42),        # Slate-900
                "muted": RGBColor(100, 116, 139),    # Slate-500
                "accent": RGBColor(59, 130, 246)     # Blue-500
            },
            "blue": {
                "bg_start": RGBColor(219, 234, 254),  # Blue-100
                "bg_end": RGBColor(239, 246, 255),   # Blue-50
                "title": RGBColor(30, 64, 175),      # Blue-700
                "body": RGBColor(15, 23, 42),
                "muted": RGBColor(96, 165, 250),     # Blue-400
                "accent": RGBColor(59, 130, 246)
            },
            "purple": {
                "bg_start": RGBColor(243, 232, 255),  # Purple-100
                "bg_end": RGBColor(250, 245, 255),    # Purple-50
                "title": RGBColor(126, 34, 206),     # Purple-700
                "body": RGBColor(30, 27, 75),
                "muted": RGBColor(168, 85, 247),    # Purple-400
                "accent": RGBColor(147, 51, 234)
            },
            "green": {
                "bg_start": RGBColor(220, 252, 231),  # Green-100
                "bg_end": RGBColor(240, 253, 244),    # Green-50
                "title": RGBColor(22, 101, 52),      # Green-700
                "body": RGBColor(20, 83, 45),
                "muted": RGBColor(74, 222, 128),     # Green-400
                "accent": RGBColor(34, 197, 94)
            },
            "orange": {
                "bg_start": RGBColor(255, 237, 213),  # Orange-100
                "bg_end": RGBColor(255, 247, 237),    # Orange-50
                "title": RGBColor(194, 65, 12),      # Orange-700
                "body": RGBColor(154, 52, 18),
                "muted": RGBColor(251, 146, 60),    # Orange-400
                "accent": RGBColor(249, 115, 22)
            },
            # --- UI presets (Corporate / Modern / Minimal) ---
            "corporate": {
                "bg_start": RGBColor(248, 250, 252),
                "bg_end": RGBColor(255, 255, 255),
                "title": RGBColor(15, 23, 42),
                "body": RGBColor(30, 41, 59),
                "muted": RGBColor(100, 116, 139),
                "accent": RGBColor(37, 99, 235),
            },
            "modern": {
                "bg_start": RGBColor(238, 242, 255),
                "bg_end": RGBColor(255, 255, 255),
                "title": RGBColor(67, 56, 202),
                "body": RGBColor(15, 23, 42),
                "muted": RGBColor(129, 140, 248),
                "accent": RGBColor(99, 102, 241),
            },
            "minimal": {
                "bg_start": RGBColor(255, 255, 255),
                "bg_end": RGBColor(250, 250, 250),
                "title": RGBColor(23, 23, 23),
                "body": RGBColor(38, 38, 38),
                "muted": RGBColor(163, 163, 163),
                "accent": RGBColor(64, 64, 64),
            },
        }
        
        # Default theme
        self.current_theme = "default"
        self._preset: str = ""

    @staticmethod
    def normalize_slide_preset(preset: Optional[str]) -> Optional[str]:
        """Trả về corporate|modern|minimal hoặc None nếu không hợp lệ / để dùng theme theo title."""
        if preset is None:
            return None
        p = str(preset).strip().lower()
        if p in VALID_SLIDE_PRESETS:
            return p
        return None

    def _clean_title(self, text: str) -> str:
        """Strip markdown heading markers (#, ##, ###) and excess whitespace."""
        t = str(text).strip()
        t = re.sub(r'^#+\s*', '', t)  # remove leading # ## ### etc.
        return t.strip()

    def _derive_chunk_title(self, bullets: List[Any], fallback: str) -> str:
        fallback_clean = self._clean_title(str(fallback or "Nội dung chính")).strip() or "Nội dung chính"
        fallback_clean = re.sub(r"\s+-\s+Ph\S*\s+\d+\s*$", "", fallback_clean, flags=re.IGNORECASE).strip() or fallback_clean
        fallback_norm = re.sub(r"\W+", " ", fallback_clean.lower()).strip()
        for raw in bullets or []:
            text = self._clean_title(str(raw or "")).strip()
            if not text:
                continue
            text = re.sub(r"^\s*(?:[-*•]|\d+[\).:-])\s*", "", text).strip()
            if ":" in text and text.find(":") <= 48:
                candidate = text.split(":", 1)[0].strip()
            else:
                first_clause = re.split(r"[.;!?]", text, maxsplit=1)[0].strip()
                comma_clause = first_clause.split(",", 1)[0].strip()
                if len(comma_clause.split()) >= 4:
                    first_clause = comma_clause
                candidate = " ".join(first_clause.split()[:11]).strip()
            candidate = self._clean_title(candidate).strip(".,;:!-“”‘’\"' ")
            candidate_norm = re.sub(r"\W+", " ", candidate.lower()).strip()
            if len(candidate.split()) >= 3 and candidate_norm and candidate_norm != fallback_norm:
                return candidate[:90]
        return fallback_clean[:90]

    @staticmethod
    def _to_inches_f(length) -> float:
        """Chuyen Length pptx hoac int EMU (sau phep tru) sang inch float."""
        if length is None:
            return 0.0
        if hasattr(length, "inches"):
            return float(length.inches)
        if isinstance(length, int):
            return length / 914400.0
        return float(length)

    def _ensure_length(self, val) -> object:
        """Chuyen int EMU (sau phep +/- Length) thanh Inches cho add_shape/textbox."""
        if val is None:
            return Inches(0)
        if hasattr(val, "inches"):
            return val
        return Inches(max(0.0, self._to_inches_f(val)))

    def _title_font_size(self, text: str, max_pt: int = 40, min_pt: int = 22) -> object:
        """Auto-scale title font size based on text length."""
        return Pt(self._title_font_points(text, max_pt=max_pt, min_pt=min_pt))

    def _title_font_points(self, text: str, max_pt: int = 40, min_pt: int = 22) -> float:
        """Kích thước font title (pt) — khớp `_title_font_size`."""
        n = len(text or "")
        if n <= 40:
            return float(max_pt)
        if n <= 60:
            return 34.0
        if n <= 80:
            return 28.0
        if n <= 110:
            return 24.0
        return float(min_pt)

    @staticmethod
    def _word_wrap_line_count(text: str, chars_per_line: float) -> int:
        """Simulate PowerPoint word-wrap: count how many lines text occupies.

        Unlike ceil(len/cpl), this respects word boundaries — a word is never
        split across lines, so the actual line count can be lower than the
        simple character-division estimate.
        """
        cpl = max(4, chars_per_line)
        words = text.split()
        if not words:
            return 1
        lines = 1
        current = 0
        for word in words:
            wlen = len(word)
            if current == 0:
                current = wlen
            elif current + 1 + wlen <= cpl:
                current += 1 + wlen
            else:
                lines += 1
                current = wlen
        return lines

    def _estimate_content_title_height_inches(
        self,
        title_text: str,
        avail_width_inches: float,
        max_pt: int = 40,
        min_pt: int = 22,
    ) -> float:
        """Ước lượng chiều cao tiêu đề bằng word-wrap simulation.

        Calibration (Inter/Calibri, 5.32-inch column with image):
          40pt → ~30 chars/line  (factor 0.32)
          34pt → ~35 chars/line  (factor 0.32, scaled)
        Verified against real slides:
          "Các yếu tố chính trong SEO"              → 1 line  ✓
          "Khái niệm và đặc điểm của Digital Marketing" → 2 lines ✓
          "Vai trò của Digital Marketing trong bối cảnh số hóa toàn cầu" → 2 lines ✓
        """
        t = (title_text or "").strip()[:200]
        if not t:
            return float(self.title_h.inches if hasattr(self.title_h, "inches") else 1.0)

        fp = self._title_font_points(t, max_pt=max_pt, min_pt=min_pt)
        # Factor 0.32: calibrated empirically for proportional Latin/Vietnamese fonts.
        # Lower than naïve 0.55 because spaces + narrow letters reduce average width.
        avg_char_in = max(0.008, (fp / 72.0) * 0.32)
        chars_per_line = avail_width_inches / avg_char_in

        # Word-wrap simulation: never splits a word → more accurate than ceil(total/cpl).
        n_lines = self._word_wrap_line_count(t, chars_per_line)
        n_lines = min(4, n_lines)

        line_h_in = (fp / 72.0) * 1.35
        pad_in = 0.06          # small padding only (box sized tightly to text)
        base = float(self.title_h.inches) if hasattr(self.title_h, "inches") else 1.0
        h = n_lines * line_h_in + pad_in
        cap = float(self.slide_height.inches) * 0.36 if hasattr(self.slide_height, "inches") else 2.5
        return float(min(cap, max(base, h)))

    def _min_title_height_for_wrap_inches(
        self,
        title_text: str,
        avail_width_inches: float,
        max_pt: int = 40,
        min_pt: int = 22,
    ) -> float:
        """Chiều cao tối thiểu theo số dòng word-wrap — không được shrink layout xuống dưới mức này."""
        t = (title_text or "").strip()[:200]
        if not t:
            return 0.81
        fp = self._title_font_points(t, max_pt=max_pt, min_pt=min_pt)
        avg_char_in = max(0.008, (fp / 72.0) * 0.32)
        chars_per_line = avail_width_inches / avg_char_in
        n_lines = min(4, self._word_wrap_line_count(t, chars_per_line))
        line_h_in = (fp / 72.0) * 1.35
        return float(n_lines * line_h_in + 0.06)

    @staticmethod
    def _safe_trunc(s: str, limit: int) -> str:
        """Cắt tại ranh giới từ hoặc câu để không cắt giữa chữ."""
        if len(s) <= limit:
            return s
        chunk = s[:limit]
        for sep in (".", "!", "?", ";", ","):
            pos = chunk.rfind(sep)
            if pos > limit // 2:
                return chunk[: pos + 1].rstrip()
        ws = chunk.rfind(" ")
        if ws > limit // 2:
            return chunk[:ws].rstrip()
        return chunk.rstrip()

    @staticmethod
    def _split_bullets_balanced(bullets: List[Any], max_bullets: int = 6) -> List[List[Any]]:
        """Split bullets into near-even chunks to avoid sparse '(tiếp)' slides.

        Example:
        - 7 bullets -> 4/3 (not 6/1)
        - 8 bullets -> 4/4
        - 10 bullets -> 5/5
        """
        items = [b for b in (bullets or []) if str(b).strip()]
        n = len(items)
        if n <= max_bullets:
            return [items] if items else []

        num_chunks = (n + max_bullets - 1) // max_bullets
        base = n // num_chunks
        rem = n % num_chunks
        chunks: List[List[Any]] = []
        start = 0
        for i in range(num_chunks):
            size = base + (1 if i < rem else 0)
            end = start + size
            chunks.append(items[start:end])
            start = end
        return chunks

    @staticmethod
    def _should_show_image(bullets: List[Any], img_path: Optional[str]) -> bool:
        """Hide image on text-sparse slides to avoid empty-looking layout."""
        if not img_path or not Path(str(img_path)).exists():
            return False
        texts = [str(b).strip() for b in (bullets or []) if str(b).strip()]
        bullet_count = len(texts)
        total_words = sum(len(re.findall(r"[\wÀ-ỹ-]+", t, flags=re.UNICODE)) for t in texts)
        if bullet_count <= 2 and total_words < 45:
            return False
        return True

    @staticmethod
    def _is_data_slide(title: str, bullets: List[Any]) -> bool:
        text = f"{title} {' '.join(str(b) for b in bullets or [])}".lower()
        has_keyword = any(kw in text for kw in _DATA_SLIDE_KEYWORDS)
        has_numbers = len(re.findall(r"[-+]?\d+(?:[.,]\d+)?\s*%?", text)) >= 2
        return has_keyword and has_numbers

    @staticmethod
    def _extract_chart_spec(title: str, bullets: List[Any]) -> Optional[Dict[str, Any]]:
        points = []
        percent_count = 0
        for idx, raw in enumerate(bullets or []):
            text = re.sub(r"\s+", " ", str(raw or "")).strip(" -•\t")
            if not text:
                continue
            match = re.search(r"[-+]?\d+(?:[.,]\d+)?\s*%?", text)
            if not match:
                continue
            raw_num = match.group(0)
            is_percent = "%" in raw_num
            number_txt = raw_num.replace("%", "").replace(",", ".").strip()
            try:
                value = float(number_txt)
            except ValueError:
                continue
            if is_percent:
                percent_count += 1
                value = value / 100.0

            label = text[: match.start()].strip(" :-–—,.;")
            if not label:
                label = text[match.end() :].strip(" :-–—,.;")
            if not label:
                label = f"Item {idx + 1}"
            points.append((label[:38], value))

        if len(points) < 2:
            return None

        labels = [p[0] for p in points[:8]]
        values = [p[1] for p in points[:8]]
        text_all = f"{title} {' '.join(str(b) for b in bullets or [])}".lower()
        if any(k in text_all for k in ("share", "composition", "distribution", "phần trăm", "phan tram", "cơ cấu", "co cau")):
            chart_type = "pie"
        elif re.search(r"\b(19\d{2}|20\d{2})\b", text_all):
            chart_type = "line"
        else:
            chart_type = "bar"
        return {
            "title": title or "Data overview",
            "labels": labels,
            "values": values,
            "chart_type": chart_type,
            "is_percent": percent_count >= max(1, len(points) // 2),
        }

    def _add_single_content_slide(
        self,
        prs,
        blank,
        title: str,
        bullets: list,
        img_path: Optional[str] = None,
        chart_spec: Optional[Dict[str, Any]] = None,
        notes: str = "",
        table_spec: Optional[Dict[str, Any]] = None,
    ) -> None:
        slide = prs.slides.add_slide(blank)
        self._add_background(slide)
        self._add_content_slide(
            slide,
            title,
            bullets,
            img_path=img_path,
            chart_spec=chart_spec,
            table_spec=table_spec,
        )
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass
    
    async def create_slide(
        self, 
        structured_content: Dict[str, Any], 
        output_path: Path,
        generate_images: bool = True,
        image_paths: Optional[Dict[int, str]] = None,
        chart_specs: Optional[Dict[int, Dict[str, Any]]] = None,
        table_specs: Optional[Dict[int, Dict[str, Any]]] = None,
        preset: Optional[str] = None,
    ):
        """
        Tạo slide PPTX
        
        Args:
            structured_content: Nội dung đã cấu trúc từ ContentExtractor
            output_path: Đường dẫn file output
            generate_images: Có sinh ảnh không
            image_paths: Dict {slide_index: image_path} nếu đã có ảnh
            chart_specs: Dict {slide_index: chart spec} cho native editable charts
            table_specs: Dict {slide_index: table spec} (headers + rows)
            preset: corporate | modern | minimal — nếu None thì giữ hành vi cũ (theme theo từ khóa title).
        """
        output_path = Path(output_path)
        
        # Tạo presentation
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        resolved = self.normalize_slide_preset(preset)
        if resolved:
            self._preset = resolved
            self.current_theme = resolved
        else:
            self._preset = ""
            theme_name = self._detect_theme(structured_content.get("title", ""))
            self.current_theme = theme_name
        
        # Always use blank layout for predictable visuals
        blank = prs.slide_layouts[6]

        # Slide tiêu đề (custom)
        title_slide = prs.slides.add_slide(blank)
        self._add_background(title_slide)
        self._add_title_slide(title_slide, structured_content.get("title", "Bài thuyết trình"))
        
        # Tạo các slide nội dung với overflow handling
        # Lọc slide trắng (không có bullet / không có bảng) trước khi render
        slides_data = [
            s for s in (structured_content.get("slides") or [])
            if s and slide_has_table_or_body(s)
        ]

        tbl_map = table_specs or {}

        for idx, slide_data in enumerate(slides_data):
            bullets = slide_data.get("bullets")
            if bullets is None:
                bullets = slide_data.get("content")  # backward compat
            bullets = [b for b in (bullets or []) if str(b).strip()]

            title_text = slide_data.get("title")
            if _is_generic_title(title_text):
                title_text = None
                if bullets:
                    first_b = str(bullets[0]).strip()
                    if ":" in first_b and first_b.find(":") < 30:
                        title_text = first_b.split(":", 1)[0].strip()
                    else:
                        words = first_b.split()
                        title_text = " ".join(words[:5]).strip('.,;:!-\u201c\u201d\u2018\u2019\'"  ')
                if not title_text or not str(title_text).strip():
                    title_text = f"Nội dung {idx + 1}"

            notes = slide_data.get("script", "") or slide_data.get("speaker_notes", "") or slide_data.get("notes", "") or ""
            table_spec_one = tbl_map.get(idx)

            # Overflow handling: split balanced to avoid sparse continuation slides.
            max_bullets = 6
            chunks = self._split_bullets_balanced(bullets, max_bullets=max_bullets)
            if not chunks:
                chunks = [[]]

            def _norm_inline_chart(raw: Optional[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
                if not raw:
                    return None
                return normalize_chart_spec(raw) if isinstance(raw, dict) else None

            if len(chunks) > 1:
                for chunk_idx, chunk in enumerate(chunks):
                    chunk_title = title_text if chunk_idx == 0 else self._derive_chunk_title(chunk, title_text)
                    base_img_path = image_paths[idx] if (generate_images and image_paths and idx in image_paths and chunk_idx == 0) else None
                    img_path = base_img_path if self._should_show_image(chunk, base_img_path) else None
                    chart_spec = None
                    if chunk_idx == 0:
                        if chart_specs and idx in chart_specs:
                            chart_spec = chart_specs[idx]
                        elif self._is_data_slide(chunk_title, chunk):
                            chart_spec = _norm_inline_chart(self._extract_chart_spec(chunk_title, chunk))
                    self._add_single_content_slide(
                        prs,
                        blank,
                        chunk_title,
                        chunk,
                        img_path=img_path,
                        chart_spec=chart_spec,
                        notes=notes if chunk_idx == 0 else "",
                        table_spec=table_spec_one if chunk_idx == 0 else None,
                    )
            else:
                base_img_path = image_paths[idx] if (generate_images and image_paths and idx in image_paths) else None
                img_path = base_img_path if self._should_show_image(bullets, base_img_path) else None
                if chart_specs and idx in chart_specs:
                    chart_spec = chart_specs[idx]
                elif self._is_data_slide(title_text, bullets):
                    chart_spec = _norm_inline_chart(self._extract_chart_spec(title_text, bullets))
                else:
                    chart_spec = None
                self._add_single_content_slide(
                    prs,
                    blank,
                    title_text,
                    bullets,
                    img_path=img_path,
                    chart_spec=chart_spec,
                    notes=notes,
                    table_spec=table_spec_one,
                )
        
        # Lưu file
        await asyncio.to_thread(prs.save, str(output_path))
        return output_path
    
    def _detect_theme(self, title: str) -> str:
        """Phát hiện theme dựa trên từ khóa trong title"""
        title_lower = title.lower()
        
        # Keywords mapping
        if any(kw in title_lower for kw in ["tech", "công nghệ", "kỹ thuật", "software", "code", "programming"]):
            return "blue"
        elif any(kw in title_lower for kw in ["marketing", "business", "kinh doanh", "bán hàng", "sales"]):
            return "orange"
        elif any(kw in title_lower for kw in ["health", "sức khỏe", "y tế", "medical", "wellness"]):
            return "green"
        elif any(kw in title_lower for kw in ["creative", "design", "thiết kế", "art", "nghệ thuật"]):
            return "purple"
        else:
            return "default"
    
    def _add_background(self, slide):
        """Thêm background với template đẹp"""
        theme = self.themes[self.current_theme]
        
        # Background chính
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.slide_width,
            self.slide_height
        )
        # Use private spTree API intentionally to enforce z-order with python-pptx.
        background.element.getparent().remove(background.element)
        slide.shapes._spTree.insert(0, background.element)
        
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = theme["bg_end"]
        background.line.fill.background()

        if self._preset == "minimal":
            top_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                self.slide_width,
                Inches(0.05),
            )
            top_line.element.getparent().remove(top_line.element)
            slide.shapes._spTree.insert(1, top_line.element)
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = theme["title"]
            top_line.line.fill.background()
            return
        
        # Thêm decorative header bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.slide_width,
            Inches(0.3)
        )
        header_bar.element.getparent().remove(header_bar.element)
        slide.shapes._spTree.insert(1, header_bar.element)
        
        header_fill = header_bar.fill
        header_fill.solid()
        header_fill.fore_color.rgb = theme["title"]
        header_bar.line.fill.background()
        
        # Thêm accent line dưới header
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(0.3),
            self.slide_width,
            Inches(0.05)
        )
        accent_line.element.getparent().remove(accent_line.element)
        slide.shapes._spTree.insert(2, accent_line.element)
        
        accent_fill = accent_line.fill
        accent_fill.solid()
        accent_fill.fore_color.rgb = theme["accent"]
        accent_line.line.fill.background()
        
    def _add_title_slide(self, slide, title: str):
        """Thêm title slide với template đẹp và dễ đọc"""
        if self._preset == "minimal":
            self._add_title_slide_minimal(slide, title)
            return

        theme = self.themes[self.current_theme]
        w = self.slide_width
        
        # Hero band for title block
        dark_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.45),
            Inches(1.55),
            w - Inches(0.9),
            Inches(3.4)
        )
        dark_fill = dark_rect.fill
        dark_fill.solid()
        dark_fill.fore_color.rgb = theme["title"]
        dark_rect.line.color.rgb = self._lighten_color(theme["accent"], 70)
        dark_rect.line.width = Pt(1.2)

        # Soft overlay for depth
        soft_overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.55),
            Inches(1.75),
            w - Inches(1.1),
            Inches(3.0)
        )
        soft_overlay.fill.solid()
        soft_overlay.fill.fore_color.rgb = self._lighten_color(theme["title"], 28)
        soft_overlay.line.fill.background()

        # Keep hero block behind text
        dark_rect.element.getparent().remove(dark_rect.element)
        slide.shapes._spTree.insert(3, dark_rect.element)
        soft_overlay.element.getparent().remove(soft_overlay.element)
        slide.shapes._spTree.insert(4, soft_overlay.element)

        # Decorative accent bars (left and right)
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8),
            Inches(2.15),
            Inches(0.11),
            Inches(2.0)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = theme["accent"]
        left_bar.line.fill.background()
        
        right_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            w - Inches(0.91),
            Inches(2.15),
            Inches(0.11),
            Inches(2.0)
        )
        right_bar.fill.solid()
        right_bar.fill.fore_color.rgb = theme["accent"]
        right_bar.line.fill.background()

        # Title box
        clean_t = self._clean_title(title)
        tb = slide.shapes.add_textbox(self.margin_x, Inches(2.3), w - 2 * self.margin_x, Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = 1
        p = tf.paragraphs[0]
        p.text = clean_t[:120]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.font_title
        p.font.size = self._title_font_size(clean_t, max_pt=50, min_pt=28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(12)

        # Decorative line under title
        title_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            w / 2 - Inches(1.8),
            Inches(4.06),
            Inches(3.6),
            Inches(0.05)
        )
        title_line_fill = title_line.fill
        title_line_fill.solid()
        title_line_fill.fore_color.rgb = self._lighten_color(theme["accent"], 85)
        title_line.line.fill.background()

        # Subtitle
        sb = slide.shapes.add_textbox(self.margin_x, Inches(4.45), w - 2 * self.margin_x, Inches(0.6))
        sf = sb.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = "Tạo bởi AI Slide Generator"
        sp.alignment = PP_ALIGN.CENTER
        sp.font.name = self.font_body
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(245, 247, 250)
        sp.font.italic = True

    def _add_title_slide_minimal(self, slide, title: str) -> None:
        """Title slide tối giản: nền sáng, chữ đậm, gạch chân mảnh."""
        theme = self.themes[self.current_theme]
        w = self.slide_width
        clean_t = self._clean_title(title)
        tb = slide.shapes.add_textbox(self.margin_x, Inches(2.45), w - 2 * self.margin_x, Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = 1
        p = tf.paragraphs[0]
        p.text = clean_t[:120]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.font_title
        p.font.size = self._title_font_size(clean_t, max_pt=44, min_pt=26)
        p.font.bold = True
        p.font.color.rgb = theme["title"]
        p.space_after = Pt(10)

        title_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            w / 2 - Inches(1.5),
            Inches(4.35),
            Inches(3.0),
            Inches(0.04),
        )
        title_line.fill.solid()
        title_line.fill.fore_color.rgb = theme["accent"]
        title_line.line.fill.background()

        sb = slide.shapes.add_textbox(self.margin_x, Inches(4.55), w - 2 * self.margin_x, Inches(0.55))
        sf = sb.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = "Tạo bởi AI Slide Generator"
        sp.alignment = PP_ALIGN.CENTER
        sp.font.name = self.font_body
        sp.font.size = Pt(15)
        sp.font.color.rgb = theme["muted"]
        sp.font.italic = True

    def _add_content_slide(
        self,
        slide,
        title: str,
        bullets: list,
        img_path: Optional[str] = None,
        chart_spec: Optional[Dict[str, Any]] = None,
        table_spec: Optional[Dict[str, Any]] = None,
    ):
        """Thêm content slide với theme colors"""
        theme = self.themes[self.current_theme]
        w = self.slide_width
        h = self.slide_height

        has_table = bool(
            table_spec
            and isinstance(table_spec.get("headers"), list)
            and table_spec.get("rows")
        )

        # Cột phải ~1/3 bề ngang vùng nội dung; trái ~2/3 (trừ khe).
        # Có bảng → không dùng cột phải (ảnh/chart) để tránh chật; ưu tiên bảng + bullet.
        has_img = bool(img_path) and Path(str(img_path)).exists() and not has_table
        has_chart = bool(chart_spec) and not has_table
        has_visual = has_img or has_chart
        inner_w_in = self._to_inches_f(w) - 2.0 * self._to_inches_f(self.margin_x)
        if has_visual:
            gap_in = 0.28
            right_col_in = inner_w_in / 3.0
            body_w_in = inner_w_in - right_col_in - gap_in
            right_col_w = Inches(right_col_in)
            gap = Inches(gap_in)
            body_w = Inches(max(3.5, body_w_in))
        else:
            right_col_w = Inches(0)
            gap = Inches(0)
            body_w = w - 2 * self.margin_x

        body_left = self.margin_x

        clean_t = self._clean_title(title)
        # Tiêu đề full chiều ngang slide (ảnh bên phải bắt đầu từ body_top, không chồng hàng title).
        # → nhiều title 1 dòng hơn, ước wrap khớp PowerPoint hơn so với cột 2/3.
        title_w = w - 2 * self.margin_x
        avail_w_in = max(4.0, self._to_inches_f(title_w))
        title_h_floor = self._min_title_height_for_wrap_inches(clean_t, avail_w_in)
        title_h_in = self._estimate_content_title_height_inches(clean_t, avail_w_in)
        title_h_in = max(title_h_in, title_h_floor)
        title_h = Inches(title_h_in)

        title_top = self.margin_y + Inches(0.05)
        # Đường accent và body bắt đầu dưới khung tiêu đề (động) — không đè lên chữ dài
        accent_top = title_top + title_h + Inches(0.06)
        body_top = accent_top + Inches(0.1)

        body_h = h - body_top - self.margin_y - self.footer_h
        min_body = Inches(1.55)
        if self._to_inches_f(body_h) < self._to_inches_f(min_body):
            # Giữ tối thiểu cho vùng bullet — KHÔNG thu title xuống dưới title_h_floor
            # (nếu không: ước 2 dòng nhưng chừa chỗ 1 dòng → khung trắng đè lên chữ).
            shrink_in = self._to_inches_f(min_body) - self._to_inches_f(body_h)
            cand = max(self._to_inches_f(self.title_h), title_h_in - shrink_in)
            title_h_in = max(title_h_floor, cand)
            title_h = Inches(title_h_in)
            accent_top = title_top + title_h + Inches(0.06)
            body_top = accent_top + Inches(0.1)
            body_h = h - body_top - self.margin_y - self.footer_h
            if self._to_inches_f(body_h) < 1.2:
                body_h = Inches(1.2)

        if not hasattr(body_h, "inches"):
            body_h = Inches(max(1.2, self._to_inches_f(body_h)))

        # Add a soft panel for reading comfort (behind text — insert trước shape chữ nếu cần)
        body_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._ensure_length(body_left - Inches(0.12)),
            self._ensure_length(body_top - Inches(0.08)),
            self._ensure_length(body_w + Inches(0.24)),
            self._ensure_length(body_h + Inches(0.05)),
        )
        body_panel.fill.solid()
        body_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        if self._preset == "minimal":
            body_panel.line.color.rgb = RGBColor(229, 229, 229)
        else:
            body_panel.line.color.rgb = self._lighten_color(theme["muted"], 90)
        body_panel.line.width = Pt(1.0)
        body_panel.element.getparent().remove(body_panel.element)
        slide.shapes._spTree.insert(3, body_panel.element)

        bullet_region_h = body_h
        table_h_in = 0.0
        if has_table:
            rows_spec = table_spec.get("rows") or []
            n_tr = min(1 + len(rows_spec), 14)
            table_h_in = min(3.75, max(0.85, 0.21 * float(n_tr) + 0.14))
            avail_in = max(0.9, self._to_inches_f(body_h) - table_h_in - 0.16)
            bullet_region_h = Inches(avail_in)

        # Title text — textbox cao hơn estimate một chút để chữ không tràn ra ngoài
        # khi word-wrap simulation sai nhẹ (edge case từ dài vừa đủ dòng).
        # Accent line vẫn dùng title_h (estimate) để định vị — không bị kéo theo buffer.
        tbox_h = self._ensure_length(title_h + Inches(0.18))
        tbox = slide.shapes.add_textbox(
            self.margin_x,
            self._ensure_length(title_top),
            self._ensure_length(title_w),
            tbox_h,
        )
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = clean_t[:120]
        p.font.name = self.font_title
        p.font.size = self._title_font_size(clean_t, max_pt=40, min_pt=22)
        p.font.bold = True
        p.font.color.rgb = theme["title"]
        p.space_after = Pt(6)

        # Accent line under title
        title_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.margin_x,
            self._ensure_length(accent_top),
            Inches(2.5),
            Inches(0.08),
        )
        title_accent_fill = title_accent.fill
        title_accent_fill.solid()
        title_accent_fill.fore_color.rgb = theme["accent"]
        title_accent.line.fill.background()

        # Body bullets
        bbox = slide.shapes.add_textbox(
            body_left,
            self._ensure_length(body_top),
            self._ensure_length(body_w),
            bullet_region_h,
        )
        btf = bbox.text_frame
        btf.word_wrap = True
        btf.clear()

        clean = []
        for b in bullets or []:
            s = str(b).strip()
            if not s:
                continue
            max_chars = 450
            if len(s) > max_chars:
                chunk = s[:max_chars].rstrip()
                last_p = max(
                    chunk.rfind("."),
                    chunk.rfind("!"),
                    chunk.rfind("?"),
                    chunk.rfind(";"),
                    chunk.rfind(":"),
                )
                if last_p > 40:
                    chunk = chunk[: last_p + 1]
                else:
                    ws = chunk.rfind(" ")
                    if ws > 40:
                        chunk = chunk[:ws].rstrip()
                s = chunk
            clean.append(s)

        if not clean and not has_table:
            clean = ["(khong co noi dung)"]

        for i, text in enumerate(clean):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            
            text_to_display = text
            font_size = Pt(19)
            if len(text) > 150:
                font_size = Pt(17)
            if len(text) > 210:
                text_to_display = self._safe_trunc(text, 210)
            
            para.text = text_to_display
            para.level = 0
            para.font.name = self.font_body
            para.font.size = font_size
            para.font.color.rgb = theme["body"]
            para.space_after = Pt(11)
            para.space_before = Pt(2)
            para.line_spacing = 1.18

        if has_table:
            t_top = self._ensure_length(body_top + bullet_region_h + Inches(0.08))
            self._add_slide_table(
                slide,
                table_spec,
                body_left,
                t_top,
                self._ensure_length(body_w),
                Inches(table_h_in),
                theme,
            )

        # Optional image/chart on right (~1/3 chiều ngang; cao gần bằng vùng body)
        if has_img:
            img_left = self._ensure_length(self.margin_x + body_w + gap)
            img_top = self._ensure_length(body_top + Inches(0.12))
            img_w = right_col_w
            bh = self._to_inches_f(body_h)
            img_h = Inches(max(2.6, min(4.35, bh - 0.1)))
            try:
                slide.shapes.add_picture(str(img_path), img_left, img_top, img_w, img_h)
            except Exception:
                pass
        elif has_chart:
            chart_left = self._ensure_length(self.margin_x + body_w + gap)
            chart_top = self._ensure_length(body_top + Inches(0.1))
            chart_w = right_col_w
            bh = self._to_inches_f(body_h)
            chart_h = Inches(max(2.7, min(4.45, bh - 0.05)))
            self._add_native_chart(slide, chart_spec or {}, chart_left, chart_top, chart_w, chart_h)

        # Footer with subtle line
        footer_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.margin_x,
            self._ensure_length(h - self.margin_y - self.footer_h - Inches(0.1)),
            self._ensure_length(w - 2 * self.margin_x),
            Inches(0.03),
        )
        footer_line_fill = footer_line.fill
        footer_line_fill.solid()
        footer_line_fill.fore_color.rgb = theme["muted"]
        footer_line.line.fill.background()
        
        fbox = slide.shapes.add_textbox(
            self.margin_x,
            self._ensure_length(h - self.margin_y - self.footer_h),
            self._ensure_length(w - 2 * self.margin_x),
            self.footer_h,
        )
        ftf = fbox.text_frame
        ftf.clear()
        fp = ftf.paragraphs[0]
        fp.text = "AI Slide Generator"
        fp.alignment = PP_ALIGN.RIGHT
        fp.font.name = self.font_body
        fp.font.size = Pt(11)
        fp.font.color.rgb = theme["muted"]

    @staticmethod
    def _xl_chart_type_for_spec(chart_kind: str, n_series: int) -> int:
        k = str(chart_kind or "bar").strip().lower().replace(" ", "_")
        if k in {"pie"}:
            return XL_CHART_TYPE.PIE
        if k in {"doughnut"}:
            return XL_CHART_TYPE.DOUGHNUT
        if k in {"line", "line_markers"}:
            return XL_CHART_TYPE.LINE_MARKERS
        if k in {"line_smooth"}:
            return XL_CHART_TYPE.LINE
        if k in {"area"}:
            return XL_CHART_TYPE.AREA
        if k in {"area_stacked"}:
            return XL_CHART_TYPE.AREA_STACKED
        if k in {"column_stacked", "col_stacked"}:
            return XL_CHART_TYPE.COLUMN_STACKED
        if k in {"column_stacked_100", "col_stacked_100"}:
            return XL_CHART_TYPE.COLUMN_STACKED_100
        if k in {"bar_horizontal", "bar_h", "hbar"}:
            return XL_CHART_TYPE.BAR_CLUSTERED
        if k in {"bar_stacked", "hbar_stacked"}:
            return XL_CHART_TYPE.BAR_STACKED
        if k in {"bar_stacked_100"}:
            return XL_CHART_TYPE.BAR_STACKED_100
        if k in {"radar"} and n_series == 1:
            return XL_CHART_TYPE.RADAR
        if k in {"bar", "column"}:
            return XL_CHART_TYPE.COLUMN_CLUSTERED
        return XL_CHART_TYPE.COLUMN_CLUSTERED

    def _add_slide_table(
        self,
        slide,
        table_spec: Dict[str, Any],
        left,
        top,
        width,
        height,
        theme,
    ) -> None:
        headers = table_spec.get("headers") or []
        rows = table_spec.get("rows") or []
        if len(headers) < 2 or not rows:
            return
        ncols = min(len(headers), 8)
        nrows = min(1 + len(rows), 14)
        try:
            shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
            tbl = shape.table
        except Exception:
            return

        for c in range(ncols):
            cell = tbl.cell(0, c)
            cell.text = str(headers[c])[:80]
            try:
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(10)
                p.font.name = self.font_body
                p.font.color.rgb = theme["title"]
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._lighten_color(theme["title"], 88)
            except Exception:
                pass

        for r in range(1, nrows):
            src = rows[r - 1] if isinstance(rows[r - 1], (list, tuple)) else []
            for c in range(ncols):
                cell = tbl.cell(r, c)
                txt = str(src[c])[:120] if c < len(src) else ""
                cell.text = txt
                try:
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(9)
                    p.font.name = self.font_body
                    p.font.color.rgb = theme["body"]
                except Exception:
                    pass

    def _add_native_chart(self, slide, chart_spec: Dict[str, Any], left, top, width, height) -> None:
        """Add an editable PowerPoint chart (multi-series + richer types)."""
        labels = [str(x)[:38] for x in (chart_spec.get("labels") or []) if str(x).strip()]
        series_list = chart_spec.get("series")
        if not isinstance(series_list, list) or not series_list:
            vals = chart_spec.get("values") or []
            if isinstance(vals, list) and labels and len(vals) >= 2:
                series_list = [
                    {
                        "name": str(chart_spec.get("title") or "Data")[:40],
                        "values": [float(v) for v in vals[: len(labels)]],
                    }
                ]
        if not labels or not series_list:
            return

        cleaned_series = []
        for s in series_list:
            if not isinstance(s, dict):
                continue
            name = str(s.get("name") or "Series")[:40]
            vs = s.get("values") or []
            if not isinstance(vs, list):
                continue
            nums = []
            for v in vs[: len(labels)]:
                try:
                    nums.append(float(v))
                except (TypeError, ValueError):
                    nums.append(0.0)
            while len(nums) < len(labels):
                nums.append(0.0)
            cleaned_series.append({"name": name, "values": nums[: len(labels)]})
        if not cleaned_series:
            return

        theme = self.themes[self.current_theme]
        chart_data = CategoryChartData()
        chart_data.categories = labels
        for s in cleaned_series:
            chart_data.add_series(s["name"], s["values"])

        n_series = len(cleaned_series)
        chart_kind = str(chart_spec.get("chart_type") or "bar").strip().lower()
        xl_type = self._xl_chart_type_for_spec(chart_kind, n_series)

        try:
            graphic_frame = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
        except Exception:
            try:
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
                )
            except Exception:
                return

        chart = graphic_frame.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = self._safe_trunc(str(chart_spec.get("title") or "Data"), 42)
        if chart_kind in {"pie", "doughnut"} or n_series <= 1:
            chart.has_legend = False
        else:
            chart.has_legend = True
            try:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            except Exception:
                pass

        if chart_kind not in {"pie", "doughnut"}:
            try:
                value_axis = chart.value_axis
                value_axis.has_major_gridlines = True
                if chart_spec.get("is_percent"):
                    value_axis.tick_labels.number_format = "0%"
            except Exception:
                pass

        try:
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(8)
            category_axis.tick_labels.font.color.rgb = theme["muted"]
        except Exception:
            pass

        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_value = True
            data_labels.font.size = Pt(8)
            if chart_spec.get("is_percent"):
                data_labels.number_format = "0%"
        except Exception:
            pass
