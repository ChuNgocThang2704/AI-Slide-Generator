import sys
import os
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

# Try finding the file in backend/outputs or root outputs
paths = [
    r"e:\DemoDoan\backend\outputs\Nội_dung_3deb1b53-9d73-4e62-b392-2d92bc1f6f4b.pptx",
    r"e:\DemoDoan\outputs\Nội_dung_3deb1b53-9d73-4e62-b392-2d92bc1f6f4b.pptx"
]

path = None
for p in paths:
    if os.path.exists(p):
        path = p
        break

if not path:
    print("Could not find the pptx file in either backend/outputs/ or outputs/")
else:
    print(f"Reading file: {path}")
    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        print(f"  Text: {text}")
