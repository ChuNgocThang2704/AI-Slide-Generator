from pptx import Presentation
import os

path = r"e:\DemoDoan\outputs\Slide_1_Tiêu_đề_và_hình_ảnh_mở_đầu_91d55afd-4fb2-4a13-b33e-ed7b04647531.pptx"
if not os.path.exists(path):
    path = r"e:\DemoDoan\outputs\Slide_1_Tiêu_đề_và_hình_ảnh_mở_đầu_ed15b07e-75e7-4ed5-a70c-b111d3f9bbd2.pptx"

output_lines = []
if not os.path.exists(path):
    output_lines.append("No such PPTX file found.")
else:
    prs = Presentation(path)
    output_lines.append(f"Reading file: {path}")
    output_lines.append(f"Total slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        output_lines.append(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        output_lines.append(f"  Text: {text}")

with open(r"e:\DemoDoan\scratch\pptx_content.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(output_lines))
print("Done writing to pptx_content.txt")
