import sys
import os
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

path = r"e:\DemoDoan\outputs\Bối_cảnh_và_tình_hình_sau_Hiệp_định_Giơ-Ne_1954_c3c89acc-0de8-4709-92f1-650cc14d5598.pptx"

if not os.path.exists(path):
    # Try backend outputs
    path = r"e:\DemoDoan\backend\outputs\Bối_cảnh_và_tình_hình_sau_Hiệp_định_Giơ-Ne_1954_c3c89acc-0de8-4709-92f1-650cc14d5598.pptx"

if not os.path.exists(path):
    print("Could not find the pptx file.")
else:
    print(f"Reading file: {path}")
    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        print(f"  Text: {text}")
